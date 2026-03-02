from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from process_reimagination_agent.models import InputManifest

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif"}
TEXT_SUFFIXES = {".txt", ".md", ".csv", ".json", ".xml"}
PDF_OCR_TEXT_THRESHOLD = 20


def _append_warning(warnings: list[str] | None, message: str) -> None:
    if warnings is not None:
        warnings.append(message)


def detect_mime_type(file_path: Path) -> str:
    try:
        import magic  # type: ignore[import-not-found]

        mime = magic.Magic(mime=True)
        return str(mime.from_file(str(file_path)))
    except Exception:
        suffix = file_path.suffix.lower()
        fallback = {
            ".pdf": "application/pdf",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".txt": "text/plain",
            ".md": "text/markdown",
            ".csv": "text/csv",
            ".json": "application/json",
            ".xml": "application/xml",
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".webp": "image/webp",
            ".bmp": "image/bmp",
            ".tiff": "image/tiff",
            ".tif": "image/tiff",
        }
        return fallback.get(suffix, "application/octet-stream")


def _ocr_image_to_text(pil_image: Any) -> str:
    try:
        import pytesseract  # type: ignore[import-not-found]
    except Exception as exc:
        raise RuntimeError("pytesseract is not installed; install it to enable OCR extraction.") from exc

    tesseract_not_found_error = getattr(pytesseract, "TesseractNotFoundError", RuntimeError)
    try:
        image = pil_image
        mode = getattr(image, "mode", "")
        if mode not in {"RGB", "L"}:
            image = image.convert("RGB")
        text = pytesseract.image_to_string(image, lang="eng")
        return str(text).strip()
    except tesseract_not_found_error as exc:  # type: ignore[misc]
        raise RuntimeError("Tesseract executable not found; install Tesseract OCR and add it to PATH.") from exc


def extract_text(file_path: Path, mime_type: str, warnings: list[str] | None = None) -> str:
    if mime_type == "application/pdf" or file_path.suffix.lower() == ".pdf":
        return extract_text_from_pdf(file_path, warnings=warnings)
    if (
        mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or file_path.suffix.lower() == ".pptx"
    ):
        return extract_text_from_pptx(file_path, warnings=warnings)
    if mime_type.startswith("image/") or file_path.suffix.lower() in IMAGE_SUFFIXES:
        return extract_text_from_image(file_path)
    if mime_type.startswith("text/") or file_path.suffix.lower() in TEXT_SUFFIXES:
        return file_path.read_text(encoding="utf-8", errors="ignore")
    return ""


def extract_text_from_image(file_path: Path) -> str:
    try:
        from PIL import Image, UnidentifiedImageError  # type: ignore[import-not-found]
    except Exception as exc:
        raise RuntimeError("Pillow is not installed; install it to enable image OCR extraction.") from exc

    try:
        with Image.open(file_path) as image:
            return _ocr_image_to_text(image)
    except UnidentifiedImageError as exc:
        raise RuntimeError(f"Could not read image file for OCR: {file_path}") from exc


def extract_text_from_pdf(file_path: Path, warnings: list[str] | None = None) -> str:
    from pypdf import PdfReader  # type: ignore[import-not-found]

    reader = PdfReader(str(file_path))
    page_texts: list[str] = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        page_texts.append(page_text.strip())

    page_images: list[Any] = []
    try:
        from pdf2image import convert_from_path  # type: ignore[import-not-found]

        page_images = list(convert_from_path(str(file_path)))
    except Exception as exc:
        _append_warning(
            warnings,
            f"{file_path}: PDF OCR rendering unavailable ({exc}). Using text-layer extraction only.",
        )

    parts: list[str] = []
    try:
        for idx, page_text in enumerate(page_texts):
            merged_text = page_text
            if len(page_text.strip()) < PDF_OCR_TEXT_THRESHOLD and idx < len(page_images):
                try:
                    ocr_text = _ocr_image_to_text(page_images[idx])
                except Exception as exc:
                    _append_warning(warnings, f"{file_path}: OCR failed on page {idx + 1}: {exc}")
                    ocr_text = ""

                if ocr_text:
                    merged_text = f"{page_text}\n{ocr_text}".strip() if page_text else ocr_text
            parts.append(merged_text)
    finally:
        for image in page_images:
            close_fn = getattr(image, "close", None)
            if callable(close_fn):
                close_fn()

    return "\n".join(part for part in parts if part).strip()


def extract_text_from_pptx(file_path: Path, warnings: list[str] | None = None) -> str:
    from pptx import Presentation  # type: ignore[import-not-found]
    from pptx.shapes.picture import Picture  # type: ignore[import-not-found]

    pil_available = True
    try:
        from PIL import Image, UnidentifiedImageError  # type: ignore[import-not-found]
    except Exception as exc:
        pil_available = False
        Image = None  # type: ignore[assignment]
        UnidentifiedImageError = Exception  # type: ignore[assignment]
        _append_warning(warnings, f"{file_path}: Pillow is not installed; skipping PPTX picture OCR ({exc}).")

    presentation = Presentation(str(file_path))
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text).strip()
                if text:
                    chunks.append(text)
            if isinstance(shape, Picture) and pil_available:
                try:
                    with Image.open(io.BytesIO(shape.image.blob)) as picture:
                        ocr_text = _ocr_image_to_text(picture)
                    if ocr_text:
                        chunks.append(ocr_text)
                except UnidentifiedImageError:
                    _append_warning(
                        warnings,
                        f"{file_path}: unsupported picture format for OCR on slide {slide_index}.",
                    )
                except Exception as exc:
                    _append_warning(
                        warnings,
                        f"{file_path}: OCR failed for picture on slide {slide_index} ({exc}).",
                    )
    return "\n".join(chunks).strip()


def ingest_manifest(manifest: InputManifest) -> dict[str, Any]:
    documents: list[dict[str, Any]] = []
    extraction_errors: list[str] = []
    combined_text_parts: list[str] = []

    for file_name in manifest.files:
        file_path = Path(file_name).expanduser()
        record: dict[str, Any] = {"path": str(file_path), "mime_type": "", "content": "", "error": ""}

        if not file_path.exists():
            record["error"] = "File does not exist"
            extraction_errors.append(f"{file_path}: not found")
            documents.append(record)
            continue

        try:
            mime_type = detect_mime_type(file_path)
            record["mime_type"] = mime_type
            extraction_warnings: list[str] = []
            content = extract_text(file_path, mime_type, warnings=extraction_warnings)
            record["content"] = content
            if extraction_warnings:
                record["warnings"] = extraction_warnings
                extraction_errors.extend(extraction_warnings)
            if content:
                combined_text_parts.append(content)
            else:
                record["error"] = "No text could be extracted"
                extraction_errors.append(f"{file_path}: empty extraction")
        except Exception as exc:
            record["error"] = f"Extraction failed: {exc}"
            extraction_errors.append(f"{file_path}: {exc}")

        documents.append(record)

    combined_text = "\n\n".join(combined_text_parts).strip()
    return {
        "manifest": manifest.model_dump(),
        "documents": documents,
        "combined_text": combined_text,
        "extraction_errors": extraction_errors,
    }
